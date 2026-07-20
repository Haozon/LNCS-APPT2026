#!/usr/bin/env python3
"""Add the 12-minute presentation script to each slide's notes (备注)."""
import os
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(HERE, "JSQKV_Defense.pptx")

NOTES = [
    # 1 — Title
    """[0:00-0:30]  Good morning / afternoon, everyone. Thank you for being here. My name is Hao Zhang, from the College of Computer Science at Nankai University. Today I'll present our work, JSQKV -- a joint sparsification and quantization framework for KV-cache compression and decode acceleration. This is joint work with my collaborators at Nankai University, the National University of Defense Technology, and the Qian Xuesen Laboratory.""",

    # 2 — Motivation
    """[0:30-1:20]  Let me start with the problem. As context windows keep growing, the real bottleneck of autoregressive decoding is no longer computation -- it's the KV cache.

The KV cache grows LINEARLY with context length, and it is re-read at EVERY single decode step. So in long-context serving, it dominates memory footprint and memory bandwidth, and it caps throughput, batch size, and deployment efficiency. The key point on the right: decoding is memory-bound -- bandwidth, not FLOPs, is the true limiter.

To address this, the community has two mainstream directions: sparsification -- evicting or pruning less useful KV states -- and low-bit quantization -- reducing the numerical precision of what's kept.""",

    # 3 — Challenge
    """[1:20-2:10]  Both directions work well on their own. The natural idea is to simply combine them -- sparsify, then quantize. But that naive chaining DOES NOT work, for two reasons.

First, a representation-and-execution gap. Sparsification changes not just how MUCH KV you keep, but the retained token and feature LAYOUT. That shifts quantization granularity, the compressed data layout, and the decode kernel -- so algorithmic compression gains don't automatically become lower memory traffic.

Second, an online-decision gap. Token-differential compression is easy in prefill, where future attention is observable. But during decoding, a token must be cached BEFORE we know how important it will be.

So the real challenge, at the bottom: preserve model quality under aggressive compression, AND keep the compressed representation directly executable for real speedup.""",

    # 4 — Overview & contributions
    """[2:10-2:55]  Our answer is JSQKV -- a unified sparse-quant decode pipeline. Instead of treating sparsity and quantization as separate post-processing steps, we co-design three things together: the compression POLICY, the compressed FORMAT, and the EXECUTION path.

Concretely, there are four contributions, shown along the bottom, and I'll walk through each: differential sparsity -- a budgeted three-level token policy; a dual-window mechanism that makes it work online; Per-Token-Tile quantization with Hadamard rotation; and a bitmap-based sparse-quant kernel that loads compressed but computes dense.""",

    # 5 — Method 1
    """[2:55-3:55]  First module: budgeted differential sparsity. The idea is simple -- not all tokens deserve the same compression. A fixed sparsity ratio wastes budget on unimportant tokens and over-compresses important ones.

So we estimate token importance from prefill attention: each token's score is the average attention it receives from the last few observation-window queries, in the SnapKV style, length-normalized.

Then we assign every token to one of three levels, shown on the right. Level 0 -- the most important -- stays dense. Level 1 keeps only the top-magnitude features in its key and value vectors. Level 2 -- the least important -- is evicted entirely. Two percentile thresholds split the tokens under a target average budget B, and a small calibration search picks the best allocation at that fixed budget. As we'll see, this matters most when the budget is tight.""",

    # 6 — Method 2
    """[3:55-4:55]  But that policy needs future attention -- which we don't have when a token is first generated. That's the second module: the dual-window mechanism.

The trick is to DELAY compression until enough future evidence accumulates. We keep two windows. Window A holds tokens waiting to be compressed. Window B holds newly generated tokens, and their queries OBSERVE Window A, supplying the missing future attention.

Three steps: we accumulate each new query's attention to Window A; we normalize by the number of observations -- since earlier tokens are seen by more queries; and we classify and slide -- comparing against the same thresholds, compressing Window A into history, and Window B becomes the new Window A.

Importantly, each token is observed before it's compressed, and the number of uncompressed tokens stays bounded. We also show these prefill thresholds are stable enough to reuse -- that's in the appendix.""",

    # 7 — Method 3
    """[4:55-5:55]  Third module: quantizing what remains. We use Per-Token-Tile quantization -- splitting each key or value vector into small tiles and quantizing each tile independently with its own scale and zero-point. This aligns the quantization unit with our token-level decisions and the tile-based decode path.

The problem is that key states carry OUTLIERS that wreck aggressive low-bit quantization. Our fix is an orthogonal Hadamard rotation applied to queries and keys before quantizing. It suppresses the heavy tails you can see in the figure -- and because it's orthogonal, it EXACTLY preserves attention scores.

The effect is dramatic. Look at the 2-bit column: plain round-to-nearest blows up to a perplexity of 115, while Hadamard rotation brings it down to 5.35 -- and it stays compatible with the sparsification that follows.""",

    # 8 — Method 4
    """[5:55-6:50]  The final module is where compression turns into REAL speedup. We co-design the storage format and the decode kernel.

Each 1-by-64 tile stores a bitmap of nonzero positions, the packed low-bit values, a per-tile offset, and the quantization metadata -- scales and zero-points.

The kernel follows a load-as-compressed, compute-as-dense strategy: it loads the compact tiles from memory, unpacks and dequantizes them into dense shared-memory tiles, and then runs a standard Tensor-Core matrix-vector product.

So sparsity is handled entirely on the LOADING path -- where the memory savings are -- while the compute stage stays dense and regular. That avoids the irregular memory access and control-flow divergence that usually kill sparse execution. This is exactly what a memory-bound decode stage needs.""",

    # 9 — Setup
    """[6:50-7:30]  Now to the evaluation. Everything runs on a single A100 80GB, with custom CUDA and Triton kernels.

For accuracy, we use the official LongBench pipeline on six representative tasks, with Llama-3-8B-Instruct as the main model, plus Llama-2, Mistral, and Qwen2.5. For efficiency, we measure end-to-end throughput and compression cost at batch sizes one through eight.

Our baselines are dense decoding, a matched-budget uniform-sparsity baseline based on MUSTAFAR, and the sequential MUSTAFAR-plus-KIVI pipeline -- which I'll call M-plus-K. We test both sparse-only and joint sparse-quant settings at 50 and 70 percent.

[PACING: if behind, trim this slide -- just name model, hardware, baselines.]""",

    # 10 — Accuracy sparse-only
    """[7:30-8:15]  First, sparsity ALONE, to isolate the effect of budget allocation -- no quantization yet.

The takeaway is on the right. Our differential policy is greater than or equal to uniform sparsity in EVERY setting -- it never hurts. And the gain is larger at the tighter 70 percent budget than at 50 percent.

That confirms the intuition from the method: token-level budget allocation matters most exactly when the budget is scarce.""",

    # 11 — Accuracy joint (main)
    """[8:15-9:15]  Now the main result -- the FULL joint sparse-quant pipeline against the sequential M-plus-K baseline, averaged over the six tasks.

JSQKV wins in ALL FOUR settings, and -- this is the key pattern -- the gap WIDENS as compression gets more aggressive. At the most aggressive setting, 70/70 with 2-bit, the average score jumps from 39.83 to 43.12 -- a gain of over three points.

Why does JSQKV win? Because jointly aligning sparsity, quantization, and execution beats stacking two methods that were never designed to work together. The cross-model results tell the same story -- favorable overall, and strongest under tight budgets.""",

    # 12 — Efficiency
    """[9:15-10:10]  And crucially, this quality holds up while delivering REAL speedup. On the left, end-to-end throughput versus batch size: the compressed paths pull ahead of dense decoding, with the clearest advantage in the small-to-medium batch regime. Overall, JSQKV improves end-to-end throughput by up to 44 percent over dense decoding.

On the right, compression statistics. Compared to sparse-only MUSTAFAR-70, our sparse-quant path pushes the compression ratio from 2.09x up to 3.48x -- and, notably, it ALSO LOWERS the compression time, from about 5,400 down to 4,300 milliseconds. So the more compact representation even cuts the write-back cost. Better compression and faster -- not a trade-off.""",

    # 13 — Conclusion
    """[10:10-11:00]  To conclude. Our central message is that KV-cache compression should be treated as a systems-and-algorithms co-design problem -- not a simple composition of independent post-processing steps.

JSQKV brings together four pieces: differential sparsity, the dual-window online mechanism, Hadamard-stabilized Per-Token-Tile quantization, and the sparse-quant kernel that turns compression into speedup.

And the headline numbers, on Llama-3-8B: the average score improves from 39.83 to 43.12, throughput improves by 44 percent, and the KV cache is compressed by 3.48 times -- all at the same time.""",

    # 14 — Thank you
    """[11:00-11:15]  That's JSQKV. Thank you very much for your attention -- I'd be happy to take your questions.

[Total ~11:15, ~45s buffer before 12:00, then 3-min Q&A.]

Likely Q&A: why differential beats uniform under tight budgets; overhead of the dual-window buffer; why gains shrink at large batch sizes (compute-bound -> dequant overhead less amortized); cross-model variance (Qwen gains big, Llama-2 modest).""",
]


def main():
    prs = Presentation(DECK)
    n = len(prs.slides._sldIdLst)
    if n != len(NOTES):
        raise SystemExit(f"Slide count {n} != notes count {len(NOTES)}; aborting.")
    for slide, text in zip(prs.slides, NOTES):
        slide.notes_slide.notes_text_frame.text = text.strip()
    prs.save(DECK)
    print(f"Added notes to {n} slides -> {DECK}")


if __name__ == "__main__":
    main()
