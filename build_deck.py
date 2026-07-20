#!/usr/bin/env python3
"""Build the JSQKV defense deck (16:9, English) with python-pptx."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = os.path.join(HERE, "ppt_assets")

# ---------- palette ----------
INK      = RGBColor(0x1F, 0x2D, 0x3D)   # near-black navy for body text
MUTED    = RGBColor(0x5B, 0x6B, 0x7B)   # secondary text
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG       = RGBColor(0xF7, 0xF9, 0xFB)   # very light page bg
CARD     = RGBColor(0xFF, 0xFF, 0xFF)

BLUE     = RGBColor(0x3D, 0x63, 0x99)   # prefill / primary brand
AMBER    = RGBColor(0xC9, 0x86, 0x2E)   # differential sparsity
GREEN    = RGBColor(0x4F, 0x8A, 0x63)   # quantization
PURPLE   = RGBColor(0x74, 0x63, 0x9E)   # dual-window
RED       = RGBColor(0xB5, 0x50, 0x50)  # kernel / highlight
DEEP     = RGBColor(0x16, 0x2A, 0x45)   # deep navy title bg

LIGHTBLUE  = RGBColor(0xE9, 0xEF, 0xF7)
LIGHTAMBER = RGBColor(0xF7, 0xEE, 0xDC)
LIGHTGREEN = RGBColor(0xE6, 0xF0, 0xE9)
LIGHTPURP  = RGBColor(0xEC, 0xE8, 0xF3)
LIGHTRED   = RGBColor(0xF6, 0xE7, 0xE7)
ROW_ALT    = RGBColor(0xF0, 0xF3, 0xF7)

FONT = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def _set_font(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def box(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         space_before=0, space_after=6, bullet=None, level=0, line=1.06, first=False):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.level = level
    try:
        p.line_spacing = line
    except Exception:
        pass
    if bullet is not None:
        r = p.add_run(); _set_font(r, size, bullet[1], bold=True); r.text = bullet[0] + "  "
    r = p.add_run(); _set_font(r, size, color, bold=bold, italic=italic); r.text = text
    return p


def rect(s, x, y, w, h, fill, line=None, line_w=0, shadow=False, round_=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if round_:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def header(s, title, accent=BLUE, kicker=None, num=None):
    """Standard content-slide header: accent bar + title + optional kicker."""
    rect(s, 0, 0, Inches(0.28), SH, accent)              # left spine
    rect(s, Inches(0.28), Inches(0.55), Inches(0.62), Inches(0.09), accent)
    tb, tf = box(s, Inches(0.55), Inches(0.66), Inches(11.6), Inches(1.0))
    if kicker:
        para(tf, kicker.upper(), 12.5, accent, bold=True, space_after=2, first=True)
        para(tf, title, 27, INK, bold=True, space_after=0)
    else:
        para(tf, title, 28, INK, bold=True, space_after=0, first=True)
    if num is not None:
        tbn, tfn = box(s, Inches(12.2), Inches(6.96), Inches(1.0), Inches(0.4))
        para(tfn, str(num), 11, MUTED, align=PP_ALIGN.RIGHT, first=True)


def pic_fit(s, path, x, y, max_w, max_h, align="center", valign="top"):
    """Add a picture scaled to fit inside (max_w,max_h) preserving aspect."""
    from PIL import Image  # noqa
    return None


def add_image(s, path, x, y, max_w, max_h, halign="center", valign="middle"):
    """Fit image within box preserving aspect ratio; returns the picture shape."""
    import struct
    # read PNG dimensions
    with open(path, "rb") as f:
        head = f.read(26)
    w_px, h_px = struct.unpack(">II", head[16:24])
    ar = w_px / h_px
    box_ar = max_w / max_h
    if ar >= box_ar:
        w = max_w; h = int(max_w / ar)
    else:
        h = max_h; w = int(max_h * ar)
    if halign == "center":
        px = x + (max_w - w) // 2
    elif halign == "right":
        px = x + (max_w - w)
    else:
        px = x
    if valign == "middle":
        py = y + (max_h - h) // 2
    elif valign == "bottom":
        py = y + (max_h - h)
    else:
        py = y
    return s.shapes.add_picture(path, px, py, width=w, height=h)


def chip(s, x, y, w, text, fill, tcolor=WHITE, h=Inches(0.42), size=12.5):
    c = rect(s, x, y, w, h, fill, round_=True)
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); _set_font(r, size, tcolor, bold=True); r.text = text
    return c


def style_table_cell(cell, text, size, color, bold=False, fill=None,
                     align=PP_ALIGN.CENTER, italic=False):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); _set_font(r, size, color, bold=bold, italic=italic)
    r.text = text


def no_table_style(tbl):
    # remove banded default style so our manual fills show
    tblPr = tbl._tbl.tblPr
    for child in list(tblPr):
        if child.tag == qn('a:tableStyleId'):
            tblPr.remove(child)
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = slide()
rect(s, 0, 0, SW, SH, DEEP)
# decorative accent stripes
rect(s, 0, Inches(6.55), SW, Inches(0.95), RGBColor(0x11, 0x20, 0x36))
for i, col in enumerate([BLUE, AMBER, GREEN, PURPLE, RED]):
    rect(s, Inches(0.0) + i * Inches(2.667), Inches(6.55), Inches(2.667), Inches(0.12), col)

tb, tf = box(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.5))
para(tf, "KV-CACHE COMPRESSION  ·  LLM INFERENCE ACCELERATION", 13,
     RGBColor(0x9F, 0xB4, 0xD4), bold=True, first=True)

tb, tf = box(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(2.2))
para(tf, "JSQKV", 60, WHITE, bold=True, space_after=6, first=True)
para(tf, "Joint Sparsification and Quantization for", 30, RGBColor(0xE3,0xEA,0xF4), bold=True, space_after=2)
para(tf, "KV-Cache Compression and Decode Acceleration", 30, RGBColor(0xE3,0xEA,0xF4), bold=True)

# authors
tb, tf = box(s, Inches(0.9), Inches(4.35), Inches(11.6), Inches(1.4))
para(tf, "Hao Zhang¹   ·   Xiaoli Gong¹*   ·   Haoran Li¹   ·   Huayou Su²   ·   Qingxia Chen³   ·   Jin Zhang¹",
     16, WHITE, bold=True, space_after=8, first=True)
para(tf, "¹ College of Computer Science, Nankai University, Tianjin, China", 12.5, RGBColor(0xB9,0xC7,0xDE), space_after=1)
para(tf, "² National University of Defense Technology, Changsha    ³ Qian Xuesen Laboratory of Space Technology, Beijing",
     12.5, RGBColor(0xB9,0xC7,0xDE), space_after=1)
para(tf, "* Corresponding author   ·   Code: github.com/Haozon/JSQKV", 11.5, RGBColor(0x8D,0xA3,0xC6), italic=True, space_before=4)

# ============================================================
# SLIDE 2 — MOTIVATION
# ============================================================
s = slide()
header(s, "The KV-Cache Is the Bottleneck of Long-Context Decoding",
       accent=BLUE, kicker="Motivation", num=2)

tb, tf = box(s, Inches(0.55), Inches(1.95), Inches(6.55), Inches(4.9))
para(tf, "Autoregressive decoding is increasingly memory-bound", 17, BLUE, bold=True,
     space_after=8, first=True)
for t in [
    "The KV cache grows linearly with context length and is re-read at every decode step.",
    "In long-context serving it dominates memory footprint and bandwidth, capping throughput, batch size, and deployment efficiency.",
    "Decoding is far more memory-bound than prefill — bandwidth, not compute, is the limiter.",
]:
    para(tf, t, 14.5, INK, bullet=("▸", BLUE), space_after=11, line=1.12)

para(tf, "Two mainstream compression directions", 17, BLUE, bold=True,
     space_before=8, space_after=8)
para(tf, "Sparsification — evict/prune less useful KV states (StreamingLLM, H₂O, SnapKV, MUSTAFAR).",
     14, INK, bullet=("•", AMBER), space_after=7, line=1.12)
para(tf, "Low-bit quantization — reduce numerical precision of retained states (KIVI, KVQuant, QuaRot).",
     14, INK, bullet=("•", GREEN), space_after=7, line=1.12)

# right stat card
cx, cy, cw = Inches(7.5), Inches(2.05), Inches(5.3)
rect(s, cx, cy, cw, Inches(4.55), CARD, line=RGBColor(0xDD,0xE3,0xEA), line_w=1, round_=True)
tb, tf = box(s, cx+Inches(0.35), cy+Inches(0.3), cw-Inches(0.7), Inches(0.6))
para(tf, "Why it matters at decode time", 14.5, INK, bold=True, first=True)
# big numbers
stats = [("Linear", "KV memory grows with sequence length", BLUE),
         ("Every step", "entire cache re-read during decoding", PURPLE),
         ("Memory-bound", "bandwidth is the true limiter, not FLOPs", RED)]
yy = cy + Inches(1.05)
for big, small, col in stats:
    rect(s, cx+Inches(0.35), yy, Inches(0.11), Inches(0.9), col)
    tb, tf = box(s, cx+Inches(0.6), yy-Inches(0.02), cw-Inches(1.0), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, big, 21, col, bold=True, space_after=1, first=True)
    para(tf, small, 12, MUTED, line=1.05)
    yy += Inches(1.12)

# ============================================================
# SLIDE 3 — CHALLENGE / GAP
# ============================================================
s = slide()
header(s, "Naïve “Sparsify-then-Quantize” Does Not Work",
       accent=RED, kicker="Problem & Challenge", num=3)

tb, tf = box(s, Inches(0.55), Inches(1.9), Inches(12.3), Inches(0.7))
para(tf, "Both directions are effective, but they are still studied in isolation. Simply chaining them fails to "
         "jointly preserve quality and deliver real end-to-end speedup.", 14.5, INK, first=True, line=1.15)

# two challenge cards
c1x, c2x = Inches(0.55), Inches(6.75)
cw2 = Inches(6.05); cy = Inches(3.0); ch = Inches(2.55)
rect(s, c1x, cy, cw2, ch, LIGHTRED, line=RED, line_w=1.25, round_=True)
tb, tf = box(s, c1x+Inches(0.35), cy+Inches(0.28), cw2-Inches(0.7), ch-Inches(0.5))
para(tf, "①  Representation & execution gap", 16.5, RED, bold=True, space_after=8, first=True)
para(tf, "Sparsification changes not just how much KV is kept, but the retained token/feature layout.",
     13.5, INK, space_after=6, line=1.14)
para(tf, "This shifts quantization granularity, compressed layout, and the decode kernel — so algorithmic "
         "compression gains do not automatically become lower memory traffic.", 13.5, INK, line=1.14)

rect(s, c2x, cy, cw2, ch, LIGHTPURP, line=PURPLE, line_w=1.25, round_=True)
tb, tf = box(s, c2x+Inches(0.35), cy+Inches(0.28), cw2-Inches(0.7), ch-Inches(0.5))
para(tf, "②  Online decision gap", 16.5, PURPLE, bold=True, space_after=8, first=True)
para(tf, "Token-differential compression is easy to estimate in prefill, where future attention is observable.",
     13.5, INK, space_after=6, line=1.14)
para(tf, "But at decode time a token must be cached before its future importance can be known — the policy "
         "is not directly applicable online.", 13.5, INK, line=1.14)

# takeaway banner
by = Inches(5.85)
rect(s, Inches(0.55), by, Inches(12.25), Inches(0.95), DEEP, round_=True)
tb, tf = box(s, Inches(0.9), by, Inches(11.6), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Key challenge: preserve model quality under aggressive compression  AND  keep the compressed "
         "representation directly executable for real decode-time acceleration.", 14.5, WHITE, bold=True,
     first=True, line=1.12)

# ============================================================
# SLIDE 4 — OVERVIEW + CONTRIBUTIONS
# ============================================================
s = slide()
header(s, "JSQKV: A Unified Sparse-Quant Decode Pipeline",
       accent=BLUE, kicker="Our Approach", num=4)
tb, tf = box(s, Inches(0.55), Inches(1.85), Inches(12.3), Inches(0.5))
para(tf, "Co-design compression policy → compressed format → execution path in one decode-stage pipeline.",
     14, MUTED, italic=True, first=True)
add_image(s, os.path.join(ASSET, "overview.png"),
          Inches(0.55), Inches(2.35), Inches(12.25), Inches(3.05), valign="top")

# 4 contribution chips
labels = [("① Differential Sparsity", "budgeted 3-level token policy", AMBER, LIGHTAMBER),
          ("② Dual-Window Online", "apply policy during decoding", PURPLE, LIGHTPURP),
          ("③ Per-Token-Tile Quant", "Hadamard-stabilized low-bit", GREEN, LIGHTGREEN),
          ("④ Sparse-Quant Kernel", "load-compressed, compute-dense", RED, LIGHTRED)]
x = Inches(0.55); w = Inches(2.96); gap = Inches(0.135); y = Inches(5.7)
for title, sub, col, lc in labels:
    rect(s, x, y, w, Inches(1.15), lc, line=col, line_w=1.25, round_=True)
    rect(s, x, y, w, Inches(0.1), col)
    tb, tf = box(s, x+Inches(0.18), y+Inches(0.2), w-Inches(0.36), Inches(0.9))
    para(tf, title, 13.5, col, bold=True, space_after=3, first=True, line=1.05)
    para(tf, sub, 11.5, INK, line=1.05)
    x += w + gap

# ============================================================
# SLIDE 5 — METHOD 1: DIFFERENTIAL SPARSITY
# ============================================================
s = slide()
header(s, "Budgeted Differential Sparsity Policy",
       accent=AMBER, kicker="Method · Module 1", num=5)

tb, tf = box(s, Inches(0.55), Inches(1.9), Inches(6.4), Inches(4.9))
para(tf, "Idea: not all tokens deserve the same compression.", 16, AMBER, bold=True,
     space_after=6, first=True)
para(tf, "A fixed sparsity ratio wastes budget on unimportant tokens and over-compresses important ones.",
     13.5, INK, space_after=12, line=1.14)

para(tf, "Estimate importance from prefill attention", 14.5, INK, bold=True, space_after=5)
para(tf, "Score = mean attention a token receives from the last L_obs observation-window queries (SnapKV-style), "
         "length-normalized.", 13, INK, bullet=("▸", AMBER), space_after=10, line=1.14)

para(tf, "Assign each token to one of three levels", 14.5, INK, bold=True, space_after=6)
para(tf, "Two percentile thresholds (τ_high, τ_low) split tokens under the target average budget B = p₁ρ₁ + p₂.",
     13, INK, bullet=("▸", AMBER), space_after=6, line=1.14)
para(tf, "A small calibration search picks the (p₀, ρ₁) allocation at fixed B.",
     13, INK, bullet=("▸", AMBER), line=1.14)

# right: three level cards
rx = Inches(7.35); rw = Inches(5.45)
levels = [("Level 0 — Dense", "Most important tokens kept fully (ρ₀ = 0).", GREEN, LIGHTGREEN),
          ("Level 1 — Partial Sparse", "Keep only top-magnitude features per key/value vector (ratio ρ₁).", AMBER, LIGHTAMBER),
          ("Level 2 — Evict", "Least important tokens removed from cache (ρ₂ = 1).", RED, LIGHTRED)]
yy = Inches(2.0)
for title, sub, col, lc in levels:
    ch = Inches(1.28)
    rect(s, rx, yy, rw, ch, lc, line=col, line_w=1.25, round_=True)
    rect(s, rx, yy, Inches(0.12), ch, col)
    tb, tf = box(s, rx+Inches(0.35), yy+Inches(0.18), rw-Inches(0.6), ch-Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 15.5, col, bold=True, space_after=4, first=True)
    para(tf, sub, 12.5, INK, line=1.1)
    yy += ch + Inches(0.16)
tb, tf = box(s, rx, yy, rw, Inches(0.5))
para(tf, "Result: larger gains at tighter (70%) budgets — token-level allocation matters most when budget is scarce.",
     12, MUTED, italic=True, first=True, line=1.1)

# ============================================================
# SLIDE 6 — METHOD 2: DUAL WINDOW
# ============================================================
s = slide()
header(s, "Dual-Window Online Execution Mechanism",
       accent=PURPLE, kicker="Method · Module 2", num=6)
tb, tf = box(s, Inches(0.55), Inches(1.9), Inches(12.3), Inches(0.75))
para(tf, "Challenge: the differential policy needs future attention, which is unavailable when a token is first "
         "generated. Solution: delay compression until enough future evidence is accumulated.", 14.5, INK,
     first=True, line=1.15)

# window diagram
dy = Inches(2.95); dh = Inches(1.35)
def window(x, w, fill, line, label, sub):
    rect(s, x, dy, w, dh, fill, line=line, line_w=1.5, round_=True)
    tb, tf = box(s, x, dy+Inches(0.22), w, dh-Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, 15, line, bold=True, align=PP_ALIGN.CENTER, space_after=3, first=True)
    para(tf, sub, 12, INK, align=PP_ALIGN.CENTER, line=1.05)
window(Inches(0.55), Inches(3.0), RGBColor(0xE4,0xE8,0xEE), MUTED, "Compressed History", "already fixed")
window(Inches(3.75), Inches(4.2), RGBColor(0xF3,0xDF,0xDF), RED, "Window A  (to compress)", "accumulates future attention")
window(Inches(8.15), Inches(4.65), RGBColor(0xDD,0xE6,0xF4), BLUE, "Window B  (observation)", "newly generated tokens supply evidence")
# arrow
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.05), Inches(4.55), Inches(2.4), Inches(0.14))
ar.fill.solid(); ar.fill.fore_color.rgb = PURPLE; ar.line.fill.background(); ar.shadow.inherit=False
tb, tf = box(s, Inches(5.6), Inches(4.75), Inches(3.3), Inches(0.3))
para(tf, "queries in B observe A", 11.5, PURPLE, italic=True, align=PP_ALIGN.CENTER, first=True)

# 3 steps
sy = Inches(5.25)
steps = [("1 · Accumulate", "Each new query's attention to Window A is summed in a buffer while B fills.", PURPLE),
         ("2 · Normalize", "Divide by #observations N⁽ⁱ⁾ = 2W−1−i, since earlier tokens are seen by more queries.", PURPLE),
         ("3 · Classify & slide", "Compare to reused (τ_high, τ_low) → merge compressed A; B becomes new A.", PURPLE)]
x = Inches(0.55); w = Inches(3.99); gap = Inches(0.14)
for t, sub, col in steps:
    rect(s, x, sy, w, Inches(1.35), CARD, line=RGBColor(0xDD,0xD8,0xE8), line_w=1.25, round_=True)
    rect(s, x, sy, w, Inches(0.08), col)
    tb, tf = box(s, x+Inches(0.22), sy+Inches(0.2), w-Inches(0.44), Inches(1.05))
    para(tf, t, 14.5, col, bold=True, space_after=4, first=True)
    para(tf, sub, 12, INK, line=1.12)
    x += w + gap
tb, tf = box(s, Inches(0.55), Inches(6.78), Inches(12), Inches(0.4))
para(tf, "Each token is observed before compression, while uncompressed tokens stay bounded by O(W). "
         "Prefill-estimated thresholds are shown to be stable enough to reuse (Appendix).",
     11.5, MUTED, italic=True, first=True)

# ============================================================
# SLIDE 7 — METHOD 3: QUANTIZATION + HADAMARD
# ============================================================
s = slide()
header(s, "Per-Token-Tile Quantization with Hadamard Rotation",
       accent=GREEN, kicker="Method · Module 3", num=7)

tb, tf = box(s, Inches(0.55), Inches(1.9), Inches(6.0), Inches(3.2))
para(tf, "Align the quantization unit with token-level decisions.", 15.5, GREEN, bold=True,
     space_after=7, first=True)
para(tf, "Split each key/value vector into contiguous tiles (size g); quantize each tile independently "
         "with its own scale & zero-point.", 13.5, INK, bullet=("▸", GREEN), space_after=8, line=1.14)
para(tf, "Key states carry outliers that destabilize low-bit quantization.", 13.5, INK,
     bullet=("▸", GREEN), space_after=8, line=1.14)
para(tf, "Apply an orthogonal Hadamard rotation to Q and K before quantizing — it suppresses heavy tails "
         "while exactly preserving attention scores (q̃ᵀk̃ = qᵀk).", 13.5, INK, bullet=("▸", GREEN), line=1.14)

# hadamard figure
add_image(s, os.path.join(ASSET, "hadamard.png"),
          Inches(6.7), Inches(1.95), Inches(6.1), Inches(2.05), valign="top")
tb, tf = box(s, Inches(6.7), Inches(3.95), Inches(6.1), Inches(0.4))
para(tf, "Key-state distribution before/after rotation (Llama-2-7B): heavy tail suppressed.",
     11, MUTED, italic=True, align=PP_ALIGN.CENTER, first=True)

# PPL table
tb, tf = box(s, Inches(0.55), Inches(4.75), Inches(12.3), Inches(0.4))
para(tf, "WikiText-2 PPL — RTN vs. Hadamard rotation (lower is better)", 13.5, INK, bold=True, first=True)

rows, cols = 3, 7
tx, ty = Inches(0.55), Inches(5.2)
tw, th = Inches(7.9), Inches(1.55)
gtbl = s.shapes.add_table(rows, cols, tx, ty, tw, th).table
no_table_style(gtbl)
gtbl.columns[0].width = Inches(2.5)
for i in range(1, 7):
    gtbl.columns[i].width = Inches((7.9-2.5)/6)
hdr = ["Setting", "4-bit RTN", "4-bit Rot.", "3-bit RTN", "3-bit Rot.", "2-bit RTN", "2-bit Rot."]
data = [
    ["Quant only", "5.22", "5.14", "6.57", "5.28", "115.55", "5.35"],
    ["+50% KV sparsity", "5.27", "5.69", "6.57", "6.35", "115.55", "6.65"],
]
for j, htext in enumerate(hdr):
    style_table_cell(gtbl.cell(0, j), htext, 10.5, WHITE, bold=True, fill=GREEN,
                     align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
for i, row in enumerate(data, start=1):
    for j, val in enumerate(row):
        best = (j in (2,4,6))  # rotate cols
        col = GREEN if best and val in ("5.14","5.28","5.35","6.35","6.65") else INK
        style_table_cell(gtbl.cell(i, j), val, 10.5,
                         col, bold=best and val in ("5.14","5.28","5.35","6.35","6.65"),
                         fill=WHITE if i % 2 else ROW_ALT,
                         align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

# highlight card
hx = Inches(8.75)
rect(s, hx, ty, Inches(4.05), th, LIGHTGREEN, line=GREEN, line_w=1.25, round_=True)
tb, tf = box(s, hx+Inches(0.28), ty+Inches(0.16), Inches(3.5), th-Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "115.55 → 5.35", 21, GREEN, bold=True, space_after=3, first=True)
para(tf, "2-bit PPL collapse is eliminated by Hadamard rotation, while staying compatible with "
         "subsequent sparsification.", 12, INK, line=1.12)

# ============================================================
# SLIDE 8 — METHOD 4: FORMAT + KERNEL
# ============================================================
s = slide()
header(s, "Bitmap-Based Sparse-Quant Format & Decode Kernel",
       accent=RED, kicker="Method · Module 4", num=8)
tb, tf = box(s, Inches(0.55), Inches(1.88), Inches(12.3), Inches(0.7))
para(tf, "Co-design storage format and kernel so compression cuts decode-time memory traffic — not just "
         "storage footprint.", 14.5, INK, first=True, line=1.15)

add_image(s, os.path.join(ASSET, "sq_format.png"),
          Inches(0.55), Inches(2.7), Inches(5.9), Inches(2.75), valign="top")
add_image(s, os.path.join(ASSET, "sq_operator.png"),
          Inches(6.75), Inches(2.55), Inches(3.35), Inches(2.95), valign="top")

tb, tf = box(s, Inches(0.55), Inches(5.5), Inches(5.9), Inches(0.35))
para(tf, "(a) Format: bitmap + packed low-bit values + metadata", 11.5, MUTED, italic=True,
     align=PP_ALIGN.CENTER, first=True)
tb, tf = box(s, Inches(6.6), Inches(5.5), Inches(3.6), Inches(0.35))
para(tf, "(b) Load-as-compressed, compute-as-dense", 11.5, MUTED, italic=True,
     align=PP_ALIGN.CENTER, first=True)

# right explanation card
ex = Inches(10.35); ew = Inches(2.45)
rect(s, ex, Inches(2.55), ew, Inches(2.95), DEEP, round_=True)
tb, tf = box(s, ex+Inches(0.22), Inches(2.75), ew-Inches(0.44), Inches(2.6))
para(tf, "Each 1×64 tile stores:", 12.5, WHITE, bold=True, space_after=5, first=True)
for t in ["Bitmap of nonzeros", "Packed low-bit values", "Per-tile offset", "Scales / zero-points"]:
    para(tf, t, 11.5, RGBColor(0xCF,0xDA,0xEC), bullet=("·", RED), space_after=4, line=1.05)
para(tf, "Kernel loads compressed tiles, unpacks + dequantizes into dense shared memory, then runs "
         "Tensor-Core MatVec.", 11, RGBColor(0xCF,0xDA,0xEC), space_before=6, line=1.12)

tb, tf = box(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(1.0))
para(tf, "Sparsity is handled on the loading path; the compute stage stays dense and regular — avoiding "
         "irregular access and control-flow divergence that usually limit sparse execution.", 13, INK,
     first=True, line=1.15, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9 — EXPERIMENTAL SETUP
# ============================================================
s = slide()
header(s, "Experimental Setup", accent=BLUE, kicker="Evaluation", num=9)

cards = [
    ("Hardware & Stack", BLUE, LIGHTBLUE,
     ["1 × NVIDIA A100 80GB PCIe", "Intel Xeon CPU, 252 GB RAM",
      "Python 3.10 · PyTorch 2.6 · CUDA 12.4", "Custom CUDA / Triton kernels"]),
    ("Models", GREEN, LIGHTGREEN,
     ["Llama-3-8B-Instruct  (main)", "Llama-2-7B", "Mistral-7B-v0.1", "Qwen2.5-7B-Instruct"]),
    ("Accuracy Protocol", AMBER, LIGHTAMBER,
     ["LongBench official pipeline", "6 tasks: NarrativeQA, Qasper,", "MultiFieldQA-EN, HotpotQA, TREC, LCC",
      "Inputs truncated to 8192 (head-tail)"]),
    ("Efficiency Protocol", RED, LIGHTRED,
     ["Llama-3-8B, 4096 in / 256 out", "Batch size 1 – 8", "Throughput, compression time,", "KV size & ratio"]),
]
x = Inches(0.55); w = Inches(2.98); gap = Inches(0.135); y = Inches(2.05); h = Inches(2.75)
for title, col, lc, items in cards:
    rect(s, x, y, w, h, lc, line=col, line_w=1.25, round_=True)
    rect(s, x, y, w, Inches(0.5), col, round_=True)
    rect(s, x, y+Inches(0.25), w, Inches(0.25), col)
    tb, tf = box(s, x+Inches(0.12), y+Inches(0.06), w-Inches(0.24), Inches(0.42), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 13.5, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    tb, tf = box(s, x+Inches(0.2), y+Inches(0.68), w-Inches(0.4), h-Inches(0.8))
    for it in items:
        para(tf, it, 11.8, INK, bullet=("·", col), space_after=6, line=1.08,
             first=(it == items[0]))
    x += w + gap

# baselines row
by = Inches(5.15)
rect(s, Inches(0.55), by, Inches(12.25), Inches(1.75), CARD, line=RGBColor(0xDD,0xE3,0xEA), line_w=1, round_=True)
tb, tf = box(s, Inches(0.85), by+Inches(0.18), Inches(11.7), Inches(0.4))
para(tf, "Baselines & compression settings", 14.5, INK, bold=True, first=True)
tb, tf = box(s, Inches(0.85), by+Inches(0.62), Inches(11.7), Inches(1.05))
para(tf, "Baselines:  Dense decoding   ·   matched-budget Uniform sparsity (MUSTAFAR)   ·   sequential "
         "MUSTAFAR + KIVI  (denoted M+K).", 13, INK, bullet=("▸", BLUE), space_after=7, first=True, line=1.12)
para(tf, "Settings:  sparse-only at 50% / 70%   ·   joint sparse-quant at 50/50 & 70/70 K/V sparsity with "
         "4-bit and 2-bit quantization. Differential policies selected on a calibration split, re-checked on a "
         "disjoint validation split.", 13, INK, bullet=("▸", AMBER), line=1.12)

# ============================================================
# SLIDE 10 — ACCURACY: SPARSE-ONLY
# ============================================================
s = slide()
header(s, "Accuracy — Differential Sparsity Alone",
       accent=AMBER, kicker="Results · 1/3", num=10)
tb, tf = box(s, Inches(0.55), Inches(1.9), Inches(12.3), Inches(0.55))
para(tf, "Isolating budget allocation (no quantization): Llama-3-8B-Instruct, 6 LongBench tasks, matched budget.",
     14, MUTED, italic=True, first=True)

# table
rows, cols = 6, 8
tx, ty = Inches(0.55), Inches(2.6)
tw, th = Inches(8.55), Inches(3.35)
t = s.shapes.add_table(rows, cols, tx, ty, tw, th).table
no_table_style(t)
t.columns[0].width = Inches(2.05)
for i in range(1, 8):
    t.columns[i].width = Inches((8.55-2.05)/7)
hdr = ["Method", "NrtvQA", "Qasper", "MF-QA", "HotpotQA", "TREC", "LCC", "Avg."]
rowsdata = [
    ("Dense", ["23.45","42.18","43.14","46.06","74.00","57.12","47.66"], "ref"),
    ("Uniform-50", ["23.44","43.73","42.93","45.94","73.50","56.03","47.59"], "base"),
    ("Differential-50", ["23.44","43.73","42.93","45.94","74.00","57.08","47.85"], "ours"),
    ("Uniform-70", ["23.94","40.89","40.91","44.93","70.00","54.12","45.80"], "base"),
    ("Differential-70", ["23.94","41.03","41.61","45.12","72.50","55.45","46.61"], "ours"),
]
for j, h in enumerate(hdr):
    style_table_cell(t.cell(0, j), h, 11, WHITE, bold=True, fill=AMBER,
                     align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
for i, (name, vals, kind) in enumerate(rowsdata, start=1):
    if kind == "ours":
        rfill = LIGHTAMBER; ncol = AMBER; nb = True
    elif kind == "ref":
        rfill = RGBColor(0xEC,0xEF,0xF3); ncol = MUTED; nb = True
    else:
        rfill = WHITE if i % 2 else ROW_ALT; ncol = INK; nb = False
    style_table_cell(t.cell(i, 0), name, 11, ncol, bold=nb, fill=rfill, align=PP_ALIGN.LEFT)
    for j, v in enumerate(vals, start=1):
        avgcol = AMBER if (kind == "ours" and j == 7) else (ncol if j == 7 else INK)
        style_table_cell(t.cell(i, j), v, 11, avgcol,
                         bold=(kind == "ours" and j == 7), fill=rfill)

# takeaway card
rx = Inches(9.4); rw = Inches(3.4)
rect(s, rx, ty, rw, th, LIGHTAMBER, line=AMBER, line_w=1.25, round_=True)
tb, tf = box(s, rx+Inches(0.28), ty+Inches(0.28), rw-Inches(0.56), th-Inches(0.5))
para(tf, "Takeaways", 15, AMBER, bold=True, space_after=8, first=True)
para(tf, "Differential ≥ uniform in every setting — never negative.", 13, INK, bullet=("▸",AMBER), space_after=9, line=1.14)
para(tf, "+0.81 avg at 70% vs +0.26 at 50%.", 13, INK, bullet=("▸",AMBER), space_after=9, line=1.14)
para(tf, "Token-level allocation matters more as the budget tightens.", 13, INK, bullet=("▸",AMBER), line=1.14)

# ============================================================
# SLIDE 11 — ACCURACY: JOINT SPARSE-QUANT (MAIN)
# ============================================================
s = slide()
header(s, "Accuracy — Joint Sparse-Quant (Main Result)",
       accent=GREEN, kicker="Results · 2/3", num=11)
tb, tf = box(s, Inches(0.55), Inches(1.9), Inches(12.3), Inches(0.5))
para(tf, "Llama-3-8B-Instruct vs. sequential MUSTAFAR+KIVI (M+K) at matched budgets — average over 6 tasks.",
     14, MUTED, italic=True, first=True)

rows, cols = 5, 3
tx, ty = Inches(0.55), Inches(2.55)
tw, th = Inches(6.5), Inches(3.5)
t = s.shapes.add_table(rows, cols, tx, ty, tw, th).table
no_table_style(t)
t.columns[0].width = Inches(2.9); t.columns[1].width = Inches(1.8); t.columns[2].width = Inches(1.8)
style_table_cell(t.cell(0,0), "Setting", 12.5, WHITE, bold=True, fill=GREEN, align=PP_ALIGN.LEFT)
style_table_cell(t.cell(0,1), "M+K", 12.5, WHITE, bold=True, fill=GREEN)
style_table_cell(t.cell(0,2), "JSQKV", 12.5, WHITE, bold=True, fill=GREEN)
jdata = [
    ("50/50 + 4-bit", "45.90", "45.93", "+0.03"),
    ("50/50 + 2-bit", "43.38", "44.02", "+0.64"),
    ("70/70 + 4-bit", "44.18", "45.34", "+1.16"),
    ("70/70 + 2-bit", "39.83", "43.12", "+3.29"),
]
for i, (name, mk, ours, dlt) in enumerate(jdata, start=1):
    hi = (i == 4)
    rfill = LIGHTGREEN if hi else (WHITE if i % 2 else ROW_ALT)
    style_table_cell(t.cell(i,0), name, 12.5, INK, bold=hi, fill=rfill, align=PP_ALIGN.LEFT)
    style_table_cell(t.cell(i,1), mk, 12.5, MUTED, fill=rfill)
    style_table_cell(t.cell(i,2), ours + f"   ({dlt})", 12.5, GREEN, bold=True, fill=rfill)

# big highlight on right
rx = Inches(7.35); rw = Inches(5.45)
rect(s, rx, ty, rw, Inches(1.65), DEEP, round_=True)
tb, tf = box(s, rx+Inches(0.35), ty+Inches(0.22), rw-Inches(0.7), Inches(1.3), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "39.83  →  43.12", 30, WHITE, bold=True, space_after=2, first=True)
para(tf, "Largest gain at the most aggressive 70/70 + 2-bit setting (+3.29 avg).", 13,
     RGBColor(0xCF,0xDA,0xEC), line=1.1)

ry2 = ty + Inches(1.85)
rect(s, rx, ry2, rw, Inches(1.65), CARD, line=RGBColor(0xDD,0xE3,0xEA), line_w=1, round_=True)
tb, tf = box(s, rx+Inches(0.35), ry2+Inches(0.2), rw-Inches(0.7), Inches(1.3))
para(tf, "Why JSQKV wins", 14.5, GREEN, bold=True, space_after=6, first=True)
para(tf, "JSQKV beats M+K in all four settings; the gap widens as compression gets more aggressive — "
         "evidence that jointly aligning sparsity, quantization and execution beats a sequential pipeline.",
     12.5, INK, line=1.14)
tb, tf = box(s, Inches(0.55), ty+th+Inches(0.05), Inches(6.5), Inches(0.4))
para(tf, "Cross-model (Llama-2 / Mistral / Qwen2.5): favorable but non-uniform, strongest under tight budgets.",
     11.5, MUTED, italic=True, first=True)

# ============================================================
# SLIDE 12 — EFFICIENCY
# ============================================================
s = slide()
header(s, "Efficiency — Throughput & Compression",
       accent=RED, kicker="Results · 3/3", num=12)

add_image(s, os.path.join(ASSET, "throughput.png"),
          Inches(0.5), Inches(2.0), Inches(6.2), Inches(3.9), valign="top")
tb, tf = box(s, Inches(0.5), Inches(5.9), Inches(6.2), Inches(0.35))
para(tf, "End-to-end decoding throughput vs. batch size (Llama-3-8B).", 11.5, MUTED, italic=True,
     align=PP_ALIGN.CENTER, first=True)

# right column: throughput stats + compression table
rx = Inches(7.0); rw = Inches(5.8)
rect(s, rx, Inches(2.0), rw, Inches(1.55), LIGHTRED, line=RED, line_w=1.25, round_=True)
tb, tf = box(s, rx+Inches(0.32), Inches(2.18), rw-Inches(0.64), Inches(1.25), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "+44.2%  end-to-end throughput", 22, RED, bold=True, space_after=3, first=True)
para(tf, "over dense decoding in long-context serving (70/70 + 2-bit); clearest gains in the "
         "small-to-medium batch regime.", 12.5, INK, line=1.12)

tb, tf = box(s, rx, Inches(3.75), rw, Inches(0.35))
para(tf, "Compression statistics (input 4096) — 70% vs. 70/70+2-bit", 13, INK, bold=True, first=True)

rows, cols = 5, 4
t = s.shapes.add_table(rows, cols, rx, Inches(4.2), rw, Inches(2.05)).table
no_table_style(t)
t.columns[0].width = Inches(2.3)
for i in range(1,4):
    t.columns[i].width = Inches((5.8-2.3)/3)
hdr = ["Method (BS=8)", "KV (MB)", "Ratio", "Time (ms)"]
cdata = [
    ("Sparse-50", "2720.2", "1.51×", "5688.2", False),
    ("Sparse-70", "1964.3", "2.09×", "5411.3", False),
    ("Sparse-50 + 2b", "1261.2", "3.25×", "4317.9", False),
    ("Sparse-70 + 2b", "1176.3", "3.48×", "4293.2", True),
]
for j,h in enumerate(hdr):
    style_table_cell(t.cell(0,j), h, 11.5, WHITE, bold=True, fill=RED,
                     align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)
for i,(name,kv,ratio,tm,hi) in enumerate(cdata, start=1):
    rfill = LIGHTRED if hi else (WHITE if i%2 else ROW_ALT)
    style_table_cell(t.cell(i,0), name, 11.5, INK, bold=hi, fill=rfill, align=PP_ALIGN.LEFT)
    style_table_cell(t.cell(i,1), kv, 11.5, INK, bold=hi, fill=rfill)
    style_table_cell(t.cell(i,2), ratio, 11.5, RED if hi else INK, bold=hi, fill=rfill)
    style_table_cell(t.cell(i,3), tm, 11.5, INK, bold=hi, fill=rfill)

tb, tf = box(s, rx, Inches(6.35), rw, Inches(0.7))
para(tf, "Higher compression ratio (2.09× → 3.48×) AND lower compression time (5411 → 4293 ms): the compact "
         "sparse-quant representation also cuts write-back cost.", 11.5, MUTED, italic=True, first=True, line=1.12)

# ============================================================
# SLIDE 13 — CONCLUSION
# ============================================================
s = slide()
header(s, "Conclusion", accent=BLUE, kicker="Summary", num=13)

tb, tf = box(s, Inches(0.55), Inches(1.95), Inches(12.2), Inches(0.85))
para(tf, "JSQKV treats KV-cache compression as a systems-and-algorithms co-design problem — not a simple "
         "composition of independent post-processing steps.", 16, INK, bold=True, first=True, line=1.18)

items = [
    ("Differential sparsity", "3-level, budgeted token allocation from prefill attention.", AMBER, LIGHTAMBER),
    ("Dual-window online", "Makes the policy applicable during autoregressive decoding.", PURPLE, LIGHTPURP),
    ("Per-Token-Tile quant", "Hadamard-stabilized low-bit representation aligned to tiles.", GREEN, LIGHTGREEN),
    ("Sparse-quant kernel", "Load-compressed / compute-dense: turns compression into speedup.", RED, LIGHTRED),
]
x = Inches(0.55); w = Inches(2.98); gap = Inches(0.135); y = Inches(3.0); h = Inches(1.75)
for title, sub, col, lc in items:
    rect(s, x, y, w, h, lc, line=col, line_w=1.25, round_=True)
    rect(s, x, y, w, Inches(0.09), col)
    tb, tf = box(s, x+Inches(0.2), y+Inches(0.24), w-Inches(0.4), h-Inches(0.4))
    para(tf, title, 14, col, bold=True, space_after=6, first=True, line=1.05)
    para(tf, sub, 12, INK, line=1.14)
    x += w + gap

# results banner
by = Inches(5.15)
rect(s, Inches(0.55), by, Inches(12.25), Inches(1.75), DEEP, round_=True)
tb, tf = box(s, Inches(0.9), by+Inches(0.22), Inches(11.6), Inches(0.5))
para(tf, "Headline results (Llama-3-8B-Instruct)", 14, RGBColor(0x9F,0xB4,0xD4), bold=True, first=True)
metrics = [("39.83 → 43.12", "avg score, 70/70+2-bit"), ("+44.2%", "end-to-end throughput"),
           ("3.48×", "KV-cache compression")]
x = Inches(1.1); w = Inches(3.7)
tby = by + Inches(0.68)
for big, small in metrics:
    tb, tf = box(s, x, tby, w, Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, big, 27, WHITE, bold=True, space_after=2, first=True)
    para(tf, small, 12.5, RGBColor(0xC5,0xD2,0xE6))
    x += Inches(3.95)

# ============================================================
# SLIDE 14 — THANK YOU / Q&A
# ============================================================
s = slide()
rect(s, 0, 0, SW, SH, DEEP)
for i, col in enumerate([BLUE, AMBER, GREEN, PURPLE, RED]):
    rect(s, i * Inches(2.667), 0, Inches(2.667), Inches(0.14), col)
    rect(s, i * Inches(2.667), Inches(7.36), Inches(2.667), Inches(0.14), col)

tb, tf = box(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(2.0), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Thank You", 54, WHITE, bold=True, space_after=10, first=True)
para(tf, "Questions & Discussion", 24, RGBColor(0xC5,0xD2,0xE6), bold=True)

tb, tf = box(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.4))
para(tf, "JSQKV — Joint Sparsification and Quantization for KV-Cache Compression and Decode Acceleration",
     14, RGBColor(0xB9,0xC7,0xDE), bold=True, first=True, space_after=6)
para(tf, "Hao Zhang, Xiaoli Gong*, Haoran Li, Huayou Su, Qingxia Chen, Jin Zhang   ·   Nankai University",
     12.5, RGBColor(0x8D,0xA3,0xC6))
para(tf, "Code:  github.com/Haozon/JSQKV", 12.5, RGBColor(0x8D,0xA3,0xC6), italic=True, space_before=3)

out = os.path.join(HERE, "JSQKV_Defense.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
